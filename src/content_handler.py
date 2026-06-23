"""Lightweight text extraction for URLs, documents, archives, and code files."""
from __future__ import annotations

import html
import logging
import os
import re
import tempfile
import zipfile
from pathlib import Path
from typing import Optional
from urllib.parse import urlparse

import httpx
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from ai import GeminiRateLimitError, gemini_youtube_video
from config import (
    MAX_ARCHIVE_FILES,
    MAX_ARCHIVE_MB,
    MAX_DOC_MB,
    MAX_TEXT_CHARS,
    MAX_URL_MB,
)

logger = logging.getLogger(__name__)

URL_RE = re.compile(r"^https?://\S+$", re.IGNORECASE)
YOUTUBE_RE = re.compile(
    r"^https?://(?:www\.)?(?:youtube\.com/watch\?v=[^&\s]+|youtu\.be/[^?\s]+|youtube\.com/shorts/[^?\s]+)",
    re.IGNORECASE,
)
CODE_EXTS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".go", ".rs", ".java", ".kt", ".c", ".h",
    ".cpp", ".hpp", ".cs", ".php", ".rb", ".swift", ".sh", ".bash", ".ps1", ".sql",
    ".json", ".jsonl", ".ndjson", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".md",
    ".txt", ".log", ".html", ".htm", ".css", ".scss", ".xml", ".dockerfile",
    ".gitignore",
}
DOC_EXTS = {".pdf", ".docx", ".pptx", ".xlsx", ".zip"}
TEXT_MIME_PREFIXES = ("text/",)
STRUCTURED_MIME = {
    "application/json",
    "application/xml",
    "text/csv",
    "application/csv",
    "application/yaml",
}


def _limit_text(text: str, limit: int = MAX_TEXT_CHARS) -> str:
    text = re.sub(r"\s+\n", "\n", text).strip()
    return text[:limit]


def _suffix_from_name(name: str | None) -> str:
    if not name:
        return ""
    low = name.lower()
    if low.endswith(".dockerfile"):
        return ".dockerfile"
    return Path(low).suffix


def _is_probably_text(data: bytes) -> bool:
    if not data:
        return True
    sample = data[:4096]
    if b"\x00" in sample:
        return False
    try:
        sample.decode("utf-8")
        return True
    except UnicodeDecodeError:
        return False


def _decode_bytes(data: bytes) -> str:
    for enc in ("utf-8", "utf-16", "cp1251", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _text_from_pdf(path: str) -> str:
    reader = PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages):
        if i >= 20:
            break
        try:
            pages.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(pages)


def _text_from_docx(path: str) -> str:
    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            vals = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if vals:
                parts.append(" | ".join(vals))
    return "\n".join(parts)


def _text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for idx, slide in enumerate(prs.slides):
        if idx >= 30:
            break
        slide_bits = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                slide_bits.append(text)
        if slide_bits:
            parts.append(f"Slide {idx + 1}:\n" + "\n".join(slide_bits))
    return "\n\n".join(parts)


def _text_from_xlsx(path: str) -> str:
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets[:5]:
        rows = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            if i >= 80:
                break
            vals = [str(v) for v in row if v is not None and str(v).strip()]
            if vals:
                rows.append(" | ".join(vals))
        if rows:
            parts.append(f"Sheet {ws.title}:\n" + "\n".join(rows))
    return "\n\n".join(parts)


def _safe_zip_names(zf: zipfile.ZipFile) -> list[str]:
    names: list[str] = []
    total_uncompressed = 0
    for info in zf.infolist():
        if len(names) >= MAX_ARCHIVE_FILES:
            break
        if info.is_dir():
            continue
        total_uncompressed += info.file_size
        if total_uncompressed > MAX_ARCHIVE_MB * 1024 * 1024:
            break
        names.append(info.filename)
    return names


def _text_from_zip(path: str) -> str:
    parts = []
    with zipfile.ZipFile(path) as zf:
        names = _safe_zip_names(zf)
        for name in names:
            if len(parts) >= 10:
                break
            suffix = _suffix_from_name(name)
            try:
                with zf.open(name) as f:
                    raw = f.read(MAX_TEXT_CHARS)
            except Exception:
                continue
            if suffix in CODE_EXTS or _is_probably_text(raw):
                text = _limit_text(_decode_bytes(raw))
                if text:
                    parts.append(f"{name}:\n{text}")
    if parts:
        return "ZIP archive text preview:\n\n" + "\n\n---\n\n".join(parts)
    with zipfile.ZipFile(path) as zf:
        names = _safe_zip_names(zf)
    return "ZIP archive with files:\n" + "\n".join(f"- {name}" for name in names)


def _text_from_plain(path: str) -> str:
    with open(path, "rb") as f:
        raw = f.read(MAX_TEXT_CHARS * 2)
    return _limit_text(_decode_bytes(raw))


def _text_from_html(html_bytes: bytes) -> str:
    soup = BeautifulSoup(html_bytes, "html.parser")
    for tag in soup(["script", "style", "noscript", "svg"]):
        tag.decompose()
    title = soup.title.get_text(" ", strip=True) if soup.title else ""
    text = soup.get_text("\n", strip=True)
    text = html.unescape(text)
    if title and title not in text[:500]:
        text = f"{title}\n\n{text}"
    return _limit_text(text)


def _is_youtube_url(url: str) -> bool:
    return bool(YOUTUBE_RE.match(url))


async def _download_url(url: str) -> tuple[str, str | None, str]:
    headers = {
        "User-Agent": "Mozilla/5.0 (CodexBot/1.0)",
        "Accept": "*/*",
    }
    max_bytes = MAX_URL_MB * 1024 * 1024
    async with httpx.AsyncClient(timeout=30, follow_redirects=True, headers=headers) as client:
        async with client.stream("GET", url) as resp:
            resp.raise_for_status()
            content_type = resp.headers.get("content-type", "").split(";", 1)[0].strip().lower() or None
            suffix = _suffix_from_name(urlparse(url).path)
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix or ".bin") as f:
                written = 0
                async for chunk in resp.aiter_bytes():
                    written += len(chunk)
                    if written > max_bytes:
                        raise ValueError(f"URL content exceeds {MAX_URL_MB} MB")
                    f.write(chunk)
                return f.name, content_type, suffix


def _is_code_or_text_name(name: str | None) -> bool:
    suffix = _suffix_from_name(name)
    return suffix in CODE_EXTS or suffix in {".csv", ".log", ".json", ".jsonl", ".ndjson", ".xml", ".yaml", ".yml"}


def _doc_kind(filename: str | None, mime_type: str | None) -> str | None:
    suffix = _suffix_from_name(filename)
    if suffix in {".pdf", ".docx", ".pptx", ".xlsx", ".zip"}:
        return suffix.lstrip(".")
    if mime_type:
        if mime_type == "application/pdf":
            return "pdf"
        if mime_type in {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"}:
            return "docx"
        if mime_type in {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}:
            return "pptx"
        if mime_type in {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"}:
            return "xlsx"
        if mime_type == "application/zip":
            return "zip"
        if mime_type.startswith(TEXT_MIME_PREFIXES) or mime_type in STRUCTURED_MIME:
            return "text"
    if suffix in CODE_EXTS:
        return "code"
    if suffix in {".txt", ".log", ".json", ".jsonl", ".ndjson", ".xml", ".csv"}:
        return "text"
    return None


def _file_kind(path: str, filename: str | None, mime_type: str | None) -> str | None:
    kind = _doc_kind(filename, mime_type)
    if kind:
        return kind
    if mime_type and mime_type.startswith(TEXT_MIME_PREFIXES):
        return "text"
    if _is_code_or_text_name(filename):
        return "code"
    with open(path, "rb") as f:
        if _is_probably_text(f.read(4096)):
            return "text"
    return None


def _analyze_local_file(path: str, filename: str | None, mime_type: str | None) -> str | None:
    kind = _file_kind(path, filename, mime_type)
    try:
        if kind == "pdf":
            return _limit_text(_text_from_pdf(path))
        if kind == "docx":
            return _limit_text(_text_from_docx(path))
        if kind == "pptx":
            return _limit_text(_text_from_pptx(path))
        if kind == "xlsx":
            return _limit_text(_text_from_xlsx(path))
        if kind == "zip":
            return _limit_text(_text_from_zip(path), MAX_TEXT_CHARS * 2)
        if kind in {"text", "code"}:
            return _limit_text(_text_from_plain(path))
    except Exception as exc:
        logger.error("Local content extraction failed: %s", exc)
        return None
    return None


async def _youtube_oembed(url: str) -> Optional[str]:
    """Fetch YouTube title + description via oEmbed as a fallback (no API key needed)."""
    try:
        async with httpx.AsyncClient(timeout=10, follow_redirects=True) as client:
            r = await client.get(
                "https://www.youtube.com/oembed",
                params={"url": url, "format": "json"},
            )
            r.raise_for_status()
            data = r.json()
            title = data.get("title", "")
            author = data.get("author_name", "")
            parts = []
            if title:
                parts.append(f"Title: {title}")
            if author:
                parts.append(f"Channel: {author}")
            return "\n".join(parts) if parts else None
    except Exception as exc:
        logger.warning("YouTube oEmbed fallback failed: %s", exc)
        return None


async def analyze_url(url: str) -> Optional[str]:
    url = url.strip()
    if not URL_RE.match(url):
        return None
    if _is_youtube_url(url):
        try:
            return await gemini_youtube_video(
                url,
                "Summarize this public YouTube video in 2-3 short sentences. Mention the main topic and any important moments. Do not provide reasoning.",
            )
        except GeminiRateLimitError:
            logger.warning("Gemini 429 — falling back to oEmbed for %s", url)
            return await _youtube_oembed(url)
        except Exception as exc:
            logger.error("YouTube Gemini analysis failed: %s — trying oEmbed", exc)
            return await _youtube_oembed(url)
    path = None
    try:
        path, content_type, suffix = await _download_url(url)
        text = _analyze_local_file(path, url, content_type)
        if text:
            return text
        if content_type and content_type.startswith("text/html"):
            return _limit_text(_text_from_html(Path(path).read_bytes()))
        if content_type and content_type.startswith(TEXT_MIME_PREFIXES):
            return _limit_text(_text_from_plain(path))
        if suffix in CODE_EXTS:
            return _limit_text(_text_from_plain(path))
        return f"Downloaded URL content type: {content_type or 'unknown'}"
    except Exception as exc:
        logger.error("URL analysis failed: %s", exc)
        return None
    finally:
        if path and os.path.exists(path):
            try:
                os.unlink(path)
            except OSError:
                pass


async def analyze_file(path: str, filename: str | None = None, mime_type: str | None = None) -> Optional[str]:
    return _analyze_local_file(path, filename, mime_type)
